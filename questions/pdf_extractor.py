import os
import PyPDF2
from pdfminer.high_level import extract_pages
from pdfminer.layout import LTTextContainer, LTChar
import pdfplumber
from PIL import Image
from pdf2image import convert_from_path
import pytesseract
from pptx import Presentation


# Function to extract text from PDF elements
def text_extraction(element):
    line_text = element.get_text()
    line_formats = []
    for text_line in element:
        if isinstance(text_line, LTTextContainer):
            for character in text_line:
                if isinstance(character, LTChar):
                    line_formats.append(character.fontname)
                    line_formats.append(character.size)
    format_per_line = list(set(line_formats))
    return (line_text, format_per_line)


# Extracting tables from the PDF page
def extract_table(pdf_path, page_num, table_num):
    pdf = pdfplumber.open(pdf_path)
    table_page = pdf.pages[page_num]
    table = table_page.extract_tables()[table_num]
    return table


# Convert the table into string format
def table_converter(table):
    table_string = ''
    for row in table:
        cleaned_row = [item.replace('\n', ' ') if item is not None and '\n' in item else 'None' if item is None else item for item in row]
        table_string += '|' + '|'.join(cleaned_row) + '|\n'
    table_string = table_string[:-1]
    return table_string


# Check if the element is inside any table
def is_element_inside_any_table(element, page, tables):
    x0, y0up, x1, y1up = element.bbox
    y0 = page.bbox[3] - y1up
    y1 = page.bbox[3] - y0up
    for table in tables:
        tx0, ty0, tx1, ty1 = table.bbox
        if tx0 <= x0 <= x1 <= tx1 and ty0 <= y0 <= y1 <= ty1:
            return True
    return False


# Find the table corresponding to an element
def find_table_for_element(element, page, tables):
    x0, y0up, x1, y1up = element.bbox
    y0 = page.bbox[3] - y1up
    y1 = page.bbox[3] - y0up
    for i, table in enumerate(tables):
        tx0, ty0, tx1, ty1 = table.bbox
        if tx0 <= x0 <= x1 <= tx1 and ty0 <= y0 <= y1 <= ty1:
            return i
    return None


# Convert the PDF to images and extract text from images
def convert_to_images(input_file):
    images = convert_from_path(input_file)
    image = images[0]
    output_file = 'PDF_image.png'
    image.save(output_file, 'PNG')


def image_to_text(image_path):
    img = Image.open(image_path)
    text = pytesseract.image_to_string(img)
    return text


# Extract text from PDF files
def extract_text_from_pdf(pdf_path):
    text_per_page = {}
    image_flag = False

    with open(pdf_path, 'rb') as pdfFileObj:
        pdfReaded = PyPDF2.PdfReader(pdfFileObj)

        for pagenum, page in enumerate(extract_pages(pdf_path)):
            pageObj = pdfReaded.pages[pagenum]
            page_text = []
            line_format = []
            text_from_images = []
            text_from_tables = []
            page_content = []
            table_in_page = -1
            pdf = pdfplumber.open(pdf_path)
            page_tables = pdf.pages[pagenum]
            tables = page_tables.find_tables()
            if len(tables) != 0:
                table_in_page = 0

            for table_num in range(len(tables)):
                table = extract_table(pdf_path, pagenum, table_num)
                table_string = table_converter(table)
                text_from_tables.append(table_string)

            page_elements = [(element.y1, element) for element in page._objs]
            page_elements.sort(key=lambda a: a[0], reverse=True)

            for i, component in enumerate(page_elements):
                element = component[1]

                if table_in_page != -1 and is_element_inside_any_table(element, page, tables):
                    table_found = find_table_for_element(element, page, tables)
                    if table_found == table_in_page and table_found is not None:
                        page_content.append(text_from_tables[table_in_page])
                        page_text.append('table')
                        line_format.append('table')
                        table_in_page += 1
                    continue

                if not is_element_inside_any_table(element, page, tables):
                    if isinstance(element, LTTextContainer):
                        (line_text, format_per_line) = text_extraction(element)
                        page_text.append(line_text)
                        line_format.append(format_per_line)
                        page_content.append(line_text)

            dctkey = 'Page_' + str(pagenum)
            text_per_page[dctkey] = [page_text, line_format, text_from_images, text_from_tables, page_content]

    full_text = ''
    for page_num, content in text_per_page.items():
        full_text += ''.join(content[0])
    return full_text


# Extract text from PPT files
def extract_text_from_ppt(ppt_path):
    presentation = Presentation(ppt_path)
    full_text = ''

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text += shape.text + '\n'
    return full_text


# Unified function to extract text based on file type
def extract_text_from_file(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()

    if file_extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_extension == '.pptx':
        return extract_text_from_ppt(file_path)
    else:
        raise ValueError("Unsupported file format. Only PDF and PPTX are supported.")

